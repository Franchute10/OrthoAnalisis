"""
OrthoAnalysis — Script de ingesta del Knowledge Base
=====================================================
Extrae texto de TXTs, PDFs digitales, DOCXs y PPTXs,
trocea en fragmentos, genera embeddings con OpenAI y los guarda en Supabase.

USO:
  python ingestar_knowledge_base.py

Variables de entorno requeridas:
  ANTHROPIC_API_KEY
  SUPABASE_URL
  SUPABASE_SERVICE_KEY
  OPENAI_API_KEY

Los archivos van en dos lugares:
  - Raíz del repo (/app/): los TXTs y DOCXs
  - Subcarpeta docs/ (/app/docs/): los PDFs y PPTXs pesados

IDEMPOTENTE: si un archivo ya fue ingestado, lo salta automáticamente.
"""

import os
import re
import time
from pathlib import Path
from pypdf import PdfReader
import openai
from supabase import create_client

# ── Configuración ─────────────────────────────────────────────────────────────

SUPABASE_URL   = os.environ.get("SUPABASE_URL")
SUPABASE_KEY   = os.environ.get("SUPABASE_SERVICE_KEY")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")

RAIZ           = Path("docs")     # todos los archivos están en docs/
CARPETA_DOCS   = Path("docs")     # todos los archivos están en docs/
TAMANO         = 400              # palabras por fragmento
SUPERPOSICION  = 50               # palabras de superposición
PAUSA          = 0.3              # segundos entre calls a OpenAI

# ── Mapa completo de archivos ─────────────────────────────────────────────────
# formato: "nombre_archivo": {"fuente": "...", "tema": "...", "carpeta": RAIZ o CARPETA_DOCS}

ARCHIVOS = {

    # ── TXTs en raíz ──────────────────────────────────────────────────────────
    "663132977-Analisis-Cefalometrico-de-Bimler-y-Sassouni.txt": {
        "fuente": "Análisis Cefalométrico de Bimler y Sassouni — R1 Ortodoncia",
        "tema": "bimler", "carpeta": RAIZ
    },
    "549077346-petrovic-clasificacion.txt": {
        "fuente": "Identificación grupo auxológico Petrovic-Lavergne — UCE Ecuador",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "631023927-petrovic.txt": {
        "fuente": "Petrovic — KIRU 2022, Vol. 19(1): 36-45",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "581357172-IDENTIFICACION-DE-TIPOS-ROTACIONALES-Y-CATEGORIAS-AUXOLOGICAS-COMO-HERRAMIENTA-DIAGNOSTICA-EN-LA-PREDICCION-DEL-POTENCIAL-DE-CRECIMIENTO-MANDIBULAR.txt": {
        "fuente": "Identificación de Tipos Rotacionales y Categorías Auxológicas — 2014",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "270275239-Identificacion-de-Tipos-Rotacionales-y-Categorias-Auxologicas-Como-Herramienta-Diagnostica-en-La-Prediccion-Del-Potencial-de-Crecimiento-Mandibular.txt": {
        "fuente": "Identificación de Tipos Rotacionales y Categorías Auxológicas — Vol 52 N°3 2014",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "384846322-La-mandibula-su-rotacion-pdf.txt": {
        "fuente": "La Mandíbula: Su Rotación — Lavergne-Petrovic",
        "tema": "rotacion_mandibular", "carpeta": RAIZ
    },
    "459980079-5-CEFALOMETRIA-DE-PETROVIC-Y-SCHWARTZ-pptx.txt": {
        "fuente": "Cefalometría de Petrovic y Schwartz — Material docente",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "507099513-Cefalometri-a-de-petrovic.txt": {
        "fuente": "Cefalometría de Petrovic — Material académico",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "295281609-analisis-cefalometrico-basico.txt": {
        "fuente": "Análisis Cefalométrico Básico — Prof. Martha Torres, UCV",
        "tema": "cefalometria", "carpeta": RAIZ
    },
    "641499029-BJORK-PREDICCION-MANDIBULAR.txt": {
        "fuente": "Björk — Predicción del Crecimiento Mandibular",
        "tema": "rotacion_mandibular", "carpeta": RAIZ
    },
    "641857580-Untitled.txt": {
        "fuente": "Material académico — Ortopedia funcional",
        "tema": "cefalometria", "carpeta": RAIZ
    },
    "585223363-AULA-PETROVIC-CEDEFACE.txt": {
        "fuente": "Aula Petrovic — CEDEFACE",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "181569359-Teoria-de-Petrovic-1.txt": {
        "fuente": "Teoría de Petrovic — Material docente",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "640479084-Crecimiento-mandibular-post-natal.txt": {
        "fuente": "Crecimiento Mandibular Postnatal",
        "tema": "rotacion_mandibular", "carpeta": RAIZ
    },
    "444359971-Sebenta-de-Ortodontia.txt": {
        "fuente": "Sebenta de Ortodontia — Universidade Católica (portugués)",
        "tema": "cefalometria", "carpeta": RAIZ
    },
    # TXTs nuevos
    "899210816-PRESENTACION-FINAL-PETROVIC-1-8-pm-1-1.txt": {
        "fuente": "Presentación Final Petrovic — Material docente",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },
    "727408514-Mandibular-Rotation-Revisited-what-makes-it-so-imp-en-es.txt": {
        "fuente": "Mandibular Rotation Revisited — traducción español",
        "tema": "rotacion_mandibular", "carpeta": RAIZ
    },
    "646789644-Petrovic-Imodelo.txt": {
        "fuente": "Petrovic — Modelo I, clasificación auxológica",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },

    # ── DOCXs en raíz (son texto plano disfrazado) ────────────────────────────
    "405807022-CEFALOMETRIA-DE-BIMLER-1-docx.docx": {
        "fuente": "Cefalometría de Bimler — Material académico",
        "tema": "bimler", "carpeta": RAIZ
    },
    "270275239-Identificacion-de-Tipos-Rotacionales-y-Categorias-Auxologicas-Como-Herramienta-Diagnostica-en-La-Prediccion-Del-Potencial-de-Crecimiento-Mandibular.docx": {
        "fuente": "Identificación de Tipos Rotacionales — Material académico",
        "tema": "tipos_rotacionales", "carpeta": RAIZ
    },

    # ── PDFs con texto en raíz ────────────────────────────────────────────────
    "480465385CefalometriaBimler.pdf": {
        "fuente": "Cefalometría de Bimler — Universidad Antonio Nariño 2020",
        "tema": "bimler", "carpeta": RAIZ
    },
    "463586775APARATOLOGIADESIMOES.pdf": {
        "fuente": "Aparatología de Simões — Ortopedia funcional",
        "tema": "aparatologia", "carpeta": RAIZ
    },
    "LUIZA_DINIZ.pdf": {
        "fuente": "Importância da Ortopedia Funcional dos Maxilares — Luíza Diniz 2020",
        "tema": "ortopedia_funcional", "carpeta": RAIZ
    },

    # ── Archivos en docs/ ─────────────────────────────────────────────────────
    "690063858-Analisis-Cefalometrico-de-Bimler-PASO-1.pdf": {
        "fuente": "Análisis Cefalométrico de Bimler — Paso a Paso",
        "tema": "bimler", "carpeta": CARPETA_DOCS
    },
    "970716306-AULA-PETROVIC-CEDEFACE.pdf": {
        "fuente": "Aula Petrovic — CEDEFACE (PDF completo)",
        "tema": "tipos_rotacionales", "carpeta": CARPETA_DOCS
    },
    "400904610-La-utilizacion-de-aparatos-funcionales-y-posturales-en-o.pdf": {
        "fuente": "Utilización de aparatos funcionales y posturales en ortopedia",
        "tema": "aparatologia", "carpeta": CARPETA_DOCS
    },
    "cefalometria-de-petrovic-y-schwartz-pptx-230604020535-38802e38.pdf": {
        "fuente": "Cefalometría de Petrovic y Schwartz — PPT completo",
        "tema": "tipos_rotacionales", "carpeta": CARPETA_DOCS
    },
    "975311262-The-Bimler.pdf": {
        "fuente": "The Bimler — Material en inglés",
        "tema": "bimler", "carpeta": CARPETA_DOCS
    },
    "A_study_to_establish_a_formula_for_predicting_the_.pdf": {
        "fuente": "A study to establish a formula for predicting mandibular growth",
        "tema": "rotacion_mandibular", "carpeta": CARPETA_DOCS
    },
    "385913751-Cefalometria-I-Bimler-Ccc.doc": {
        "fuente": "Cefalometría I — Bimler CCC",
        "tema": "bimler", "carpeta": CARPETA_DOCS
    },
    "532123603-Bimler.pptx": {
        "fuente": "Bimler — Presentación académica completa",
        "tema": "bimler", "carpeta": CARPETA_DOCS
    },
}

# ── Extracción de texto ───────────────────────────────────────────────────────

def extraer_txt(ruta):
    with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extraer_pdf(ruta):
    reader = PdfReader(str(ruta))
    texto = ""
    for page in reader.pages:
        texto += (page.extract_text() or "") + "\n"
    return texto

def extraer_docx_plano(ruta):
    """DOCXs descargados de Scribd son texto plano con extensión .docx"""
    with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extraer_pptx(ruta):
    """Extrae texto de PPTX usando python-pptx si está disponible."""
    try:
        from pptx import Presentation
        prs = Presentation(str(ruta))
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto
    except ImportError:
        print("  ⚠️  python-pptx no disponible — saltando PPTX")
        return ""
    except Exception as e:
        print(f"  ⚠️  Error PPTX: {e}")
        return ""

def extraer_texto(ruta):
    ext = ruta.suffix.lower()
    if ext == '.txt':
        return extraer_txt(ruta)
    elif ext == '.pdf':
        return extraer_pdf(ruta)
    elif ext in ('.docx', '.doc'):
        return extraer_docx_plano(ruta)
    elif ext == '.pptx':
        return extraer_pptx(ruta)
    return ""

# ── Fragmentación ─────────────────────────────────────────────────────────────

def limpiar(texto):
    texto = re.sub(r'\n{3,}', '\n\n', texto)
    texto = re.sub(r' {2,}', ' ', texto)
    texto = re.sub(r'\f', '\n\n', texto)
    return texto.strip()

def trocear(texto, tamano=TAMANO, superposicion=SUPERPOSICION):
    palabras = texto.split()
    fragmentos = []
    i = 0
    while i < len(palabras):
        fin = min(i + tamano, len(palabras))
        fragmento = " ".join(palabras[i:fin])
        if len(fragmento.strip()) > 100:
            fragmentos.append(fragmento)
        i += tamano - superposicion
    return fragmentos

# ── Embeddings con OpenAI ─────────────────────────────────────────────────────

def generar_embedding(cliente_openai, texto):
    for intento in range(3):
        try:
            response = cliente_openai.embeddings.create(
                input=texto,
                model="text-embedding-3-small"  # 1536 dims, $0.02/1M tokens
            )
            return response.data[0].embedding
        except Exception as e:
            if intento < 2:
                print(f"\n  ⏳ Rate limit, reintentando ({intento+1}/3)...")
                time.sleep(2 ** intento)
            else:
                raise e

# ── Supabase ──────────────────────────────────────────────────────────────────

def ya_ingestado(sb, nombre):
    result = sb.table("knowledge_base").select("id").eq("documento", nombre).limit(1).execute()
    return len(result.data) > 0

def guardar(sb, texto, embedding, fuente, documento, fragmento_n, tema):
    sb.table("knowledge_base").insert({
        "texto":     texto,
        "fuente":    fuente,
        "documento": documento,
        "pagina":    fragmento_n,
        "tema":      tema,
        "acceso":    "global",
        "embedding": embedding
    }).execute()

# ── Proceso principal ─────────────────────────────────────────────────────────

def procesar(nombre, config, sb, openai_client):
    ruta = config["carpeta"] / nombre

    if not ruta.exists():
        print(f"  ⚠️  No encontrado: {ruta}")
        return 0

    if ya_ingestado(sb, nombre):
        print(f"  ⏭️  Ya ingestado")
        return 0

    print(f"  📄 Extrayendo texto...")
    texto = limpiar(extraer_texto(ruta))

    if len(texto) < 200:
        print(f"  ⚠️  Texto muy corto ({len(texto)} chars) — saltando")
        return 0

    fragmentos = trocear(texto)
    print(f"  ✂️  {len(fragmentos)} fragmentos de {TAMANO} palabras")

    guardados = 0
    for i, frag in enumerate(fragmentos):
        try:
            print(f"  🔢 Embedding {i+1}/{len(fragmentos)}...", end="\r")
            emb = generar_embedding(openai_client, frag)
            guardar(sb, frag, emb, config["fuente"], nombre, i+1, config["tema"])
            guardados += 1
            time.sleep(PAUSA)
        except Exception as e:
            print(f"\n  ❌ Error fragmento {i+1}: {e}")
            continue

    print(f"\n  ✅ {guardados}/{len(fragmentos)} fragmentos guardados")
    return guardados

def main():
    print("=" * 60)
    print("OrthoAnalysis — Ingesta Knowledge Base")
    print("=" * 60)

    # Verificar variables
    errores = []
    if not SUPABASE_URL:   errores.append("Falta SUPABASE_URL")
    if not SUPABASE_KEY:   errores.append("Falta SUPABASE_SERVICE_KEY")
    if not OPENAI_API_KEY: errores.append("Falta OPENAI_API_KEY")
    if errores:
        print("\n❌ Variables faltantes:")
        for e in errores: print(f"   - {e}")
        return

    # Conectar
    print("\n🔌 Conectando servicios...")
    sb = create_client(SUPABASE_URL, SUPABASE_KEY)
    oa = openai.OpenAI(api_key=OPENAI_API_KEY)
    print("✅ Supabase y OpenAI conectados")

    # Crear carpeta docs si no existe
    CARPETA_DOCS.mkdir(exist_ok=True)

    # Procesar
    total = 0
    print(f"\n📚 Procesando {len(ARCHIVOS)} archivos...\n")

    for nombre, config in ARCHIVOS.items():
        print(f"\n{'─' * 50}")
        print(f"📖 {nombre[:55]}")
        print(f"   Tema: {config['tema']} | {config['fuente'][:45]}")
        n = procesar(nombre, config, sb, oa)
        total += n

    print(f"\n{'=' * 60}")
    print(f"✅ INGESTA COMPLETA — {total} fragmentos guardados en knowledge_base")
    print(f"{'=' * 60}")

if __name__ == "__main__":
    main()
